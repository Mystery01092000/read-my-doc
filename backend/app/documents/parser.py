"""File-type dispatch for text extraction from uploaded documents."""

import csv
import io
from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class ParsedPage:
    page_number: int | None
    section_heading: str | None
    text: str


def parse_file(path: Path, file_type: str) -> list[ParsedPage]:
    """Extract text from a file, returning a list of pages/sections."""
    dispatch = {
        "pdf": _parse_pdf,
        "txt": _parse_text,
        "md": _parse_markdown,
        "csv": _parse_csv,
        "xlsx": _parse_excel,
        "pptx": _parse_pptx,
    }
    parser = dispatch.get(file_type)
    if parser is None:
        raise ValueError(f"Unsupported file type: {file_type}")
    return parser(path)


def _parse_pdf(path: Path) -> list[ParsedPage]:
    import fitz  # pymupdf

    pages: list[ParsedPage] = []
    with fitz.open(str(path)) as doc:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append(ParsedPage(page_number=page_num, section_heading=None, text=text))
    return pages


def _parse_text(path: Path) -> list[ParsedPage]:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    if not text:
        return []
    return [ParsedPage(page_number=None, section_heading=None, text=text)]


def _parse_markdown(path: Path) -> list[ParsedPage]:
    """Split markdown by top-level headings, preserving section context."""
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    sections: list[ParsedPage] = []
    current_heading: str | None = None
    current_lines: list[str] = []

    for line in lines:
        if line.startswith("# "):
            if current_lines:
                content = "\n".join(current_lines).strip()
                if content:
                    sections.append(ParsedPage(None, current_heading, content))
            current_heading = line[2:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        content = "\n".join(current_lines).strip()
        if content:
            sections.append(ParsedPage(None, current_heading, content))

    if not sections:
        return [ParsedPage(None, None, text.strip())]
    return sections


def _parse_csv(path: Path) -> list[ParsedPage]:
    """Convert CSV to readable text (header + rows)."""
    text_parts: list[str] = []
    with path.open(encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        return []

    headers = rows[0]
    chunk_size = 50  # rows per chunk
    for i in range(1, len(rows), chunk_size):
        batch = rows[i : i + chunk_size]
        lines = [", ".join(f"{h}: {v}" for h, v in zip(headers, row)) for row in batch]
        text_parts.append("\n".join(lines))

    return [
        ParsedPage(page_number=idx + 1, section_heading=None, text=chunk)
        for idx, chunk in enumerate(text_parts)
    ]


def _parse_excel(path: Path) -> list[ParsedPage]:
    """Convert Excel sheets to text."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    pages: list[ParsedPage] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(c) if c is not None else "" for c in rows[0]]
        text_lines: list[str] = []
        for row in rows[1:]:
            line = ", ".join(
                f"{h}: {v}" for h, v in zip(headers, row) if v is not None and str(v).strip()
            )
            if line:
                text_lines.append(line)

        if text_lines:
            pages.append(ParsedPage(page_number=None, section_heading=sheet_name, text="\n".join(text_lines)))

    wb.close()
    return pages


def _parse_pptx(path: Path) -> list[ParsedPage]:
    """Extract text from each PowerPoint slide."""
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[ParsedPage] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        heading: str | None = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if not line:
                    continue
                # Use the first bold/large text as slide heading
                if heading is None and para.runs and para.runs[0].font.bold:
                    heading = line
                else:
                    parts.append(line)

        text = "\n".join(parts).strip()
        if text or heading:
            pages.append(ParsedPage(page_number=slide_num, section_heading=heading, text=text or heading or ""))

    return pages
